import os
import streamlit as st
from groq import Groq
import re
import PyPDF2
import io
import base64
from docx import Document
from pptx import Presentation
import httpx  # Add httpx import
from llama_index.indices.managed.llama_cloud import LlamaCloudIndex  # For vector DB access
import streamlit.components.v1 as components
import textwrap
import ast
import time
import json
from dotenv import load_dotenv

# Load environment variables from .env file
load_dotenv()

# Check for API key in environment variables or Streamlit secrets
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
if not GROQ_API_KEY:
    st.error("⚠️ GROQ_API_KEY not found! Please set it in your .env file or environment variables.")
    st.stop()

# --- Helper for flashcard validation ---
def is_valid_flashcard(card):
    """Check if a flashcard dict has all required keys and correct types."""
    return (
        isinstance(card, dict)
        and all(k in card for k in ['question', 'options', 'correct', 'explanation'])
        and isinstance(card['options'], list)
    )

def validate_flashcard_with_details(card, index=None):
    """
    Validate a flashcard and return detailed error information.
    Returns (is_valid: bool, error_msg: str)
    """
    prefix = f"Flashcard #{index + 1}: " if index is not None else "Flashcard: "
    
    if not isinstance(card, dict):
        return False, f"{prefix}Not a dictionary object, got {type(card)}"
    
    required_keys = ['question', 'options', 'correct', 'explanation']
    missing_keys = [key for key in required_keys if key not in card]
    if missing_keys:
        return False, f"{prefix}Missing required keys: {missing_keys}"
    
    # Check question
    if not isinstance(card['question'], str) or not card['question'].strip():
        return False, f"{prefix}Question must be a non-empty string"
    
    # Check options
    if not isinstance(card['options'], list):
        return False, f"{prefix}Options must be a list, got {type(card['options'])}"
    if len(card['options']) != 4:
        return False, f"{prefix}Must have exactly 4 options, got {len(card['options'])}"
    if not all(isinstance(opt, str) and opt.strip() for opt in card['options']):
        return False, f"{prefix}All options must be non-empty strings"
    
    # Check correct answer
    if not isinstance(card['correct'], str):
        return False, f"{prefix}Correct answer must be a string, got {type(card['correct'])}"
    if card['correct'].upper() not in ['A', 'B', 'C', 'D']:
        return False, f"{prefix}Correct answer must be A, B, C, or D, got '{card['correct']}'"
    
    # Check explanation
    if not isinstance(card['explanation'], str) or not card['explanation'].strip():
        return False, f"{prefix}Explanation must be a non-empty string"
    
    return True, f"{prefix}Valid"

# --- MODE RESET HELPERS ---
def reset_flashcard_state():
    """Clear all session state related to flashcard mode."""
    for key in [
        'generated_flashcards',
        'flashcard_idx',
    ]:
        if key in st.session_state:
            del st.session_state[key]

def reset_practice_test_state():
    """Clear all session state related to practice test mode."""
    for key in [
        'questions', 'answers', 'raw_answers', 'question_types', 'options_list',
        'current_question_idx', 'user_answers', 'show_results',
        'answer_explanations', 'answer_scores', 'session_review',
    ]:
        if key in st.session_state:
            del st.session_state[key]

# Set the API key in environment variables
os.environ["GROQ_API_KEY"] = GROQ_API_KEY

# --- VECTOR DB CONFIG ---
# Load vector database configuration from environment variables
VECTOR_DB_INDEX_ID = os.environ.get("VECTOR_DB_INDEX_ID", "3d295ac4-a443-4df9-88cc-8597fb0c23a6")
VECTOR_DB_PROJECT_ID = os.environ.get("VECTOR_DB_PROJECT_ID", "874973e8-8db0-4355-993a-47c1376ab8fe") 
VECTOR_DB_ENDPOINT = f"https://api.cloud.llamaindex.ai/api/pipeline/{VECTOR_DB_INDEX_ID}/retrieve"
VECTOR_DB_ORG_ID = os.environ.get("VECTOR_DB_ORG_ID", "d5fdc60a-1af9-404c-b030-f24ff62cc04a")
# LlamaIndex API key for vector database access (optional)
VECTOR_DB_API_KEY = os.environ.get("LLAMA_INDEX_API_KEY")

# Helper to retrieve context from vector DB
@st.cache_data(show_spinner=False, ttl=600, max_entries=10)
def retrieve_vector_context(query, api_key=VECTOR_DB_API_KEY):
    """Retrieve relevant context from the LlamaIndex vector DB for a given query."""
    try:
        index = LlamaCloudIndex(
            name="class 8 Physics",  # Or make this configurable
            project_name="Default",
            organization_id=VECTOR_DB_ORG_ID,
            api_key=api_key,
        )
        nodes = index.as_retriever().retrieve(query)
        # Extract text from nodes (top 3 for brevity)
        context_chunks = []
        for node in nodes[:3]:
            # node.text or node.get("text") depending on API
            chunk = getattr(node, "text", None) or node.get("text", "")
            if chunk:
                context_chunks.append(chunk)
        return "\n\n".join(context_chunks)
    except Exception as e:
        st.warning(f"Vector DB retrieval failed: {e}")
        return ""

def parse_questions(response_text):
    """Parse the response text into separate questions with improved formatting and robust parsing."""
    questions = []
    answers = []
    question_types = []
    options_list = []
    raw_answers = []
    
    # More robust section splitting
    sections = []
    
    # Method 1: Split by ### (primary approach)
    if "###" in response_text:
        raw_sections = response_text.split("###")
        # Filter out empty sections and clean them
        for section in raw_sections:
            if section.strip():
                sections.append(section.strip())
    
    # Method 2: Manual section detection if ### splitting fails
    if not sections:
        section_patterns = [
            (r"Multiple Choice Questions?", "multiple_choice"),
            (r"Fill in the Blank Questions?", "fill_blank"), 
            (r"True (?:or|/)?False Questions?", "true_false"),
            (r"Question and Answer Questions?", "question_answer")
        ]
        
        # Find all section headers with their positions
        section_headers = []
        for pattern, q_type in section_patterns:
            for match in re.finditer(pattern, response_text, re.IGNORECASE):
                section_headers.append((match.start(), match.end(), q_type, match.group()))
        
        # Sort by position
        section_headers.sort(key=lambda x: x[0])
        
        # Extract sections based on headers
        for i, (start, end, q_type, header) in enumerate(section_headers):
            # Find the start of content (after header)
            content_start = end
            # Find the end of this section (start of next section or end of text)
            if i + 1 < len(section_headers):
                content_end = section_headers[i + 1][0]
            else:
                content_end = len(response_text)
            
            section_content = response_text[content_start:content_end].strip()
            if section_content:
                # Add header info to section for type detection
                sections.append(f"{header}\n{section_content}")
    
    # If still no sections found, treat entire text as one section
    if not sections:
        sections = [response_text]
    
    for section_idx, section in enumerate(sections):
        if not section.strip():
            continue
            
        # Determine question type from section content
        section_lower = section.lower()
        question_type = "other"
        if "multiple choice" in section_lower:
            question_type = "multiple_choice"
        elif "fill in the blank" in section_lower or "fill-in-the-blank" in section_lower:
            question_type = "fill_blank"
        elif "true" in section_lower and "false" in section_lower:
            question_type = "true_false"
        elif "question and answer" in section_lower:
            question_type = "question_answer"

        # Enhanced question extraction with multiple approaches
        question_matches = []
        answer_matches = []
        
        # Approach 1: Regex pattern matching
        # More robust pattern that handles various edge cases
        question_pattern = r'(\d+)\.\s*(.*?)(?=\n\s*Correct Answer:|$)'
        answer_pattern = r'Correct Answer:\s*(.*?)(?=\n\s*\d+\.|$)'
        
        # Find all potential question blocks first
        question_blocks = re.split(r'\n\s*(\d+)\.\s*', section)
        
        # Process question blocks
        current_questions = []
        for i in range(1, len(question_blocks), 2):  # Skip first empty element, then take pairs
            if i + 1 < len(question_blocks):
                q_num = question_blocks[i]
                q_content = question_blocks[i + 1]
                
                # Split content into question part and answer part
                if "Correct Answer:" in q_content:
                    parts = q_content.split("Correct Answer:", 1)
                    question_text = parts[0].strip()
                    answer_text = parts[1].strip() if len(parts) > 1 else ""
                    
                    # Clean up answer (remove next question start if accidentally included)
                    answer_lines = answer_text.split('\n')
                    clean_answer = answer_lines[0].strip()
                    
                    current_questions.append((q_num, question_text, clean_answer))
        
        # Fallback: Line-by-line parsing if block parsing fails
        if not current_questions:
            lines = section.split('\n')
            current_q_num = None
            current_q_content = []
            current_answer = ""
            
            i = 0
            while i < len(lines):
                line = lines[i].strip()
                if not line:
                    i += 1
                    continue
                    
                # Check if this is a new question
                q_match = re.match(r'(\d+)\.\s*(.*)', line)
                if q_match:
                    # Save previous question if exists
                    if current_q_num and current_q_content:
                        current_questions.append((current_q_num, '\n'.join(current_q_content), current_answer))
                    
                    # Start new question
                    current_q_num = q_match.group(1)
                    current_q_content = [q_match.group(2)] if q_match.group(2) else []
                    current_answer = ""
                elif line.startswith("Correct Answer:"):
                    current_answer = line.replace("Correct Answer:", "").strip()
                elif current_q_num:
                    current_q_content.append(line)
                i += 1
            
            # Don't forget the last question
            if current_q_num and current_q_content:
                current_questions.append((current_q_num, '\n'.join(current_q_content), current_answer))

        # Process each extracted question
        for q_num, q_content, answer_content in current_questions:
            try:
                # Clean up question content
                q_lines = [line.strip() for line in q_content.split('\n') if line.strip()]
                
                # Separate question text from options
                question_text = ""
                current_options = []
                
                for line in q_lines:
                    # Check for multiple choice options
                    if re.match(r'^[A-D]\)', line):
                        current_options.append(line.strip())
                    elif line and not re.match(r'^-+$', line) and not line.startswith("Correct Answer:"):
                        if question_text:
                            question_text += " "
                        question_text += line
                
                # Skip if no question text found
                if not question_text.strip():
                    continue
                
                # Clean up question text for true/false questions
                if question_type == "true_false":
                    question_text = re.sub(r'\s*\((?:True|False)\)\s*$', '', question_text, flags=re.IGNORECASE).strip()
                    question_text = re.sub(r'\s+(?:True|False)\s*$', '', question_text, flags=re.IGNORECASE).strip()
                
                # Format question for display
                formatted_question = f'<div class="question-text">{question_text}</div>'
                if question_type == "multiple_choice" and current_options:
                    formatted_question += '<div class="options-list">'
                    formatted_question += '</div>'
                
                # Process answer
                raw_answer = re.sub(r'\*\*|\*', '', answer_content).strip() if answer_content else ""
                formatted_answer = f'''<div class="answer-section">
                    <div class="answer-label">Correct Answer:</div>
                    <div class="answer-text">{answer_content}</div>
                    </div>''' if answer_content else ""
                
                # Add to results
                questions.append(formatted_question)
                answers.append(formatted_answer)
                raw_answers.append(raw_answer)
                question_types.append(question_type)
                options_list.append(current_options)
                
            except Exception as e:
                # Log parsing errors for debugging but continue
                print(f"Error parsing question {q_num}: {e}")
                continue
    
    return questions, answers, question_types, options_list, raw_answers

def extract_text_from_pdf(pdf_file):
    """Extract text content from uploaded PDF file."""
    text = ""
    try:
        # Ensure the file pointer is at the beginning
        pdf_file.seek(0)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text: # Check if text extraction was successful
                text += page_text + "\n"
        # Reset file pointer after reading if needed elsewhere, though Streamlit handles this usually
        pdf_file.seek(0)
        return text
    except Exception as e:
        st.error(f"Error processing PDF: {e}")
        return ""

def extract_text_from_docx(docx_file):
    """Extract text content from uploaded DOCX file."""
    text = ""
    try:
        # Ensure the file pointer is at the beginning
        docx_file.seek(0)
        doc = Document(docx_file)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        # Reset file pointer
        docx_file.seek(0)
        return text
    except Exception as e:
        st.error(f"Error processing DOCX: {e}")
        return ""

def extract_text_from_ppt(ppt_file):
    """Extract text content from uploaded PPT file."""
    text = ""
    try:
        # Ensure the file pointer is at the beginning
        ppt_file.seek(0)
        prs = Presentation(ppt_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
         # Reset file pointer
        ppt_file.seek(0)
        return text
    except Exception as e:
        st.error(f"Error processing PPT: {e}")
        return ""

def extract_text_from_multiple_files(files):
    """Extract text content from multiple files (PDF, DOCX, or PPT)."""
    combined_text = ""
    # Create new file objects from bytes to handle potential pointer issues
    processed_files = []
    for uploaded_file in files:
        file_bytes = uploaded_file.getvalue()
        file_like_object = io.BytesIO(file_bytes)
        # Preserve the original name and type if needed
        file_like_object.name = uploaded_file.name
        file_like_object.type = uploaded_file.type
        processed_files.append(file_like_object)


    for file in processed_files:
        try:
            # file.seek(0) # Ensure pointer is at the start for each file
            if file.name.lower().endswith('.pdf'):
                combined_text += extract_text_from_pdf(file) # Use dedicated function
            elif file.name.lower().endswith('.docx'):
                 combined_text += extract_text_from_docx(file) # Use dedicated function
            elif file.name.lower().endswith(('.ppt', '.pptx')):
                combined_text += extract_text_from_ppt(file) # Use dedicated function
        except Exception as e:
            st.error(f"Error processing {file.name}: {e}")
        # file.seek(0) # Reset pointer after processing (optional but good practice)
    return combined_text

def generate_single_question_type_batch(query, question_type, count, source_content=None, api_key=GROQ_API_KEY):
    """
    Generate a specific number of questions of a single type using JSON mode.
    This approach is more reliable for exact counts.
    """
    if count == 0:
        return []
    
    context = f"\nContext from provided documents:\n{source_content}" if source_content else ""
    
    # Define the question type specifics
    type_specs = {
        'multiple_choice': {
            'name': 'Multiple Choice',
            'format': 'question with 4 options (A, B, C, D) and correct answer',
            'example': '{"question": "What is 2+2?", "options": ["3", "4", "5", "6"], "correct": "B", "explanation": "2+2 equals 4"}'
        },
        'fill_blank': {
            'name': 'Fill in the Blank',
            'format': 'question with a blank (____) and the correct word/phrase',
            'example': '{"question": "Water boils at ____ degrees Celsius.", "correct": "100", "explanation": "Water boils at 100°C at standard pressure"}'
        },
        'true_false': {
            'name': 'True or False',
            'format': 'statement and whether it is true or false',
            'example': '{"question": "The Earth is flat.", "correct": "False", "explanation": "The Earth is spherical, not flat"}'
        },
        'question_answer': {
            'name': 'Question and Answer',
            'format': 'open-ended question with detailed answer',
            'example': '{"question": "Explain photosynthesis.", "correct": "Photosynthesis is the process by which plants convert sunlight into energy...", "explanation": "This covers the basic definition and importance of photosynthesis"}'
        }
    }
    
    spec = type_specs[question_type]
    
    # Simple, focused system prompt for JSON mode
    system_prompt = (
        f"You are an expert {spec['name']} question generator. "
        f"Generate EXACTLY {count} {spec['name'].lower()} questions in valid JSON format. "
        f"Each question should be a {spec['format']}. "
        f"Your response must be a JSON object with a 'questions' array containing exactly {count} questions. "
        f"Example format: {spec['example']}"
    )
    
    user_prompt = f"Topic: {query}{context}\n\nGenerate exactly {count} {spec['name'].lower()} questions about this topic."
    
    try:
        http_client = httpx.Client(timeout=90.0, trust_env=False)
        client = Groq(api_key=api_key, http_client=http_client)
        completion = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.1,  # Very low temperature for consistency
            top_p=0.95,
            stream=False,
            response_format={"type": "json_object"}  # JSON mode for reliability
        )
        
        result = json.loads(completion.choices[0].message.content)
        questions = result.get('questions', [])
        
        # Validate count
        if len(questions) != count:
            print(f"Warning: Expected {count} {question_type} questions, got {len(questions)}")
        
        return questions[:count]  # Ensure we don't exceed the requested count
        
    except Exception as e:
        print(f"Error generating {question_type} questions: {e}")
        return []

def generate_flashcards(query, source_content=None, num_mc=3, num_fb=3, num_tf=3, num_qa=3):
    """
    Generate flashcards using batch generation for better reliability.
    Each question type is generated separately to ensure exact counts.
    """
    # Calculate total questions
    total_questions = num_mc + num_fb + num_tf + num_qa
    
    if total_questions == 0:
        return ""
    
    # Generate each question type separately
    question_types = [
        ('multiple_choice', num_mc),
        ('fill_blank', num_fb), 
        ('true_false', num_tf),
        ('question_answer', num_qa)
    ]
    
    all_questions = []
    generated_counts = {'multiple_choice': 0, 'fill_blank': 0, 'true_false': 0, 'question_answer': 0}
    
    for i, (q_type, count) in enumerate(question_types):
        if count > 0:
            print(f"Generating {count} {q_type.replace('_', ' ')} questions...")
            batch_questions = generate_single_question_type_batch(
                query, q_type, count, source_content
            )
            generated_counts[q_type] = len(batch_questions)
            print(f"Successfully generated {len(batch_questions)}/{count} {q_type.replace('_', ' ')} questions")
            
            for q_data in batch_questions:
                # Convert JSON format to the expected text format
                formatted_q = convert_json_to_text_format(q_data, q_type)
                all_questions.append(formatted_q)
    
    # Print final summary
    total_generated = sum(generated_counts.values())
    print(f"Generation complete: {total_generated}/{total_questions} questions generated")
    print(f"Breakdown: MC={generated_counts['multiple_choice']}/{num_mc}, FB={generated_counts['fill_blank']}/{num_fb}, TF={generated_counts['true_false']}/{num_tf}, QA={generated_counts['question_answer']}/{num_qa}")
    
    # Combine all questions into the expected format
    if not all_questions:
        return ""
    
    # Organize by type for the expected output format
    output_sections = []
    section_headers = {
        'multiple_choice': '### Multiple Choice Questions',
        'fill_blank': '### Fill in the Blank Questions', 
        'true_false': '### True or False Questions',
        'question_answer': '### Question and Answer Questions'
    }
    
    current_counts = {'multiple_choice': 0, 'fill_blank': 0, 'true_false': 0, 'question_answer': 0}
    organized_questions = {'multiple_choice': [], 'fill_blank': [], 'true_false': [], 'question_answer': []}
    
    # Group questions by type (they should already be in order from batch generation)
    for q_type, count in question_types:
        if count > 0:
            start_idx = sum(c for t, c in question_types[:question_types.index((q_type, count))])
            end_idx = start_idx + count
            organized_questions[q_type] = all_questions[start_idx:end_idx]
    
    # Build the final output
    final_output = ""
    for q_type, count in question_types:
        if count > 0:
            final_output += section_headers[q_type] + "\n"
            for i, question_text in enumerate(organized_questions[q_type], 1):
                final_output += f"{i}. {question_text}\n\n"
            if q_type != question_types[-1][0]:  # Don't add extra newlines after the last section
                final_output += "\n"
    
    return final_output.strip()

def convert_json_to_text_format(q_data, q_type):
    """Convert JSON question format to the expected text format."""
    question = q_data.get('question', '')
    correct = q_data.get('correct', '')
    explanation = q_data.get('explanation', '')
    
    if q_type == 'multiple_choice':
        options = q_data.get('options', [])
        if len(options) != 4:
            options = ['Option A', 'Option B', 'Option C', 'Option D']  # Fallback
        
        formatted = f"{question}\n"
        for i, option in enumerate(options):
            formatted += f"{chr(65+i)}) {option}\n"
        formatted += f"Correct Answer: {correct}"
        
    elif q_type == 'fill_blank':
        formatted = f"{question}\nCorrect Answer: {correct}"
        
    elif q_type == 'true_false':
        formatted = f"{question}\nCorrect Answer: {correct}"
        
    elif q_type == 'question_answer':
        formatted = f"{question}\nCorrect Answer: {correct}"
    
    return formatted

def evaluate_answer_llm(question, correct_answer, user_answer, question_type, api_key=GROQ_API_KEY):
    """
    Use the LLM to evaluate the user's answer for any question type.
    Returns a tuple: (score: float [0-1], explanation: str)
    """
    # Compose the system prompt for JSON mode
    system_prompt = (
        "You are an expert grader. Evaluate the user's answer for correctness using JSON format only. "
        "Give a score from 0 (completely wrong) to 1 (fully correct). "
        "Provide a clear, concise explanation for your judgment, tailored to the question type. "
        "Respond with JSON in this exact format: {\"score\": 0.8, \"explanation\": \"Your detailed explanation here\"}"
    )
    
    user_prompt = (
        f"Question Type: {question_type}\n"
        f"Question: {question}\n"
        f"Correct Answer: {correct_answer}\n"
        f"User's Answer: {user_answer}\n\n"
        f"Evaluate the user's answer and respond with the required JSON format."
    )
    
    try:
        http_client = httpx.Client(timeout=60.0, trust_env=False)
        client = Groq(api_key=api_key, http_client=http_client)
        completion = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.0,
            top_p=0.95,
            stream=False,
            response_format={"type": "json_object"}  # JSON mode enabled!
        )
        
        # JSON mode guarantees valid JSON - no complex parsing needed!
        result = json.loads(completion.choices[0].message.content)
        score = float(result.get("score", 0))
        explanation = result.get("explanation", "No explanation provided.")
        score = max(0.0, min(1.0, score))  # Ensure score is within bounds
        return score, explanation
        
    except Exception as e:
        return 0.0, f"LLM evaluation failed: {e}"

# Initialize session state for question index and revealed answers if not exists
if 'current_question_idx' not in st.session_state:
    st.session_state.current_question_idx = 0
if 'questions' not in st.session_state:
    st.session_state.questions = []
if 'answers' not in st.session_state: # Display-formatted answers
    st.session_state.answers = []
if 'raw_answers' not in st.session_state: # Raw answers for comparison
    st.session_state.raw_answers = []
if 'revealed_answers' not in st.session_state:
    st.session_state.revealed_answers = set()
if 'user_selections' not in st.session_state:
    st.session_state.user_selections = {}
if 'question_types' not in st.session_state:
    st.session_state.question_types = []
if 'options_list' not in st.session_state:
    st.session_state.options_list = []
if 'pdf_data' not in st.session_state:
    st.session_state.pdf_data = []
if 'pdf_names' not in st.session_state:
    st.session_state.pdf_names = []
if 'user_answers' not in st.session_state:
    st.session_state.user_answers = {} # Stores answers for all types: {idx: answer}
if 'answer_feedback' not in st.session_state:
    st.session_state.answer_feedback = {}
if 'show_results' not in st.session_state:
    st.session_state.show_results = False
if 'source_content' not in st.session_state: # Store extracted text
    st.session_state.source_content = ""
if 'vector_context' not in st.session_state:
    st.session_state.vector_context = ""

# Add custom CSS for the flashcard styling
st.markdown("""
    <style>
    .flashcard {
        background-color: #ffffff;
        padding: 2rem;
        border-radius: 15px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        margin: 1rem 0 2rem 0; /* Adjusted margin */
        border: 1px solid #e0e0e0;
    }
    .stButton button {
        border-radius: 20px;
        padding: 0.5rem 1rem;
        font-weight: 500;
        margin: 0.2rem; /* Added margin for spacing */
    }
    .nav-buttons {
        display: flex;
        justify-content: space-between;
        align-items: center;
        margin-top: 1rem;
        padding: 0 1rem; /* Padding for spacing */
    }
    .question-counter {
        color: #666;
        font-size: 0.9rem;
        text-align: center;
    }
    .question-text {
        font-size: 1.1rem;
        font-weight: 500;
        margin-bottom: 1rem; /* Reduced margin */
        color: #1f1f1f;
    }
    /* Input specific styles */
    .stRadio > label { /* Target radio button labels */
        font-size: 1rem !important; /* Ensure consistent font size */
        padding: 0.5rem 0;
    }
    /* Style for single-line text input */
    .stTextInput input {
        border: 1px solid #ccc;
        border-radius: 8px;
        padding: 0.75rem;
        margin-top: none;
        width: 100%;
        box-sizing: border-box;
        box-shadow: none; /* Removed shadow */
        transition: border-color 0.2s ease; /* Removed box-shadow from transition */
        background-color: #fff; /* Ensure background is white */
    }
    /* Style for text area */
    .stTextArea textarea {
        border: 1px solid #ccc;
        border-radius: 8px;
        padding: 0.75rem;
        margin-top: none;
        width: 100%;
        box-sizing: border-box;
        box-shadow: none; /* Removed shadow */
        transition: border-color 0.2s ease; /* Removed box-shadow from transition */
        background-color: #fff; /* Ensure background is white */
    }

    /* Focus style for single-line text input */
    .stTextInput input:focus {
        border-color: #4a90e2;
        box-shadow: none; /* Removed shadow */
        outline: none;
    }
    /* Focus style for text area */
    .stTextArea textarea:focus {
        border-color: #4a90e2;
        box-shadow: none; /* Removed shadow */
        outline: none;
    }
    /* Results page styling */
    .result-item {
        border: 1px solid #ddd;
        border-radius: 10px;
        padding: 1.5rem;
        margin-bottom: 1.5rem;
        background-color: #f9f9f9;
    }
    .result-header {
        font-size: 1.1rem;
        font-weight: bold;
        margin-bottom: 1rem;
        color: #333;
    }
    .result-section {
        margin-bottom: 0.8rem;
        padding-left: 1rem;
        border-left: 3px solid #eee;
    }
    .result-label {
        font-weight: bold;
        color: #555;
        display: block; /* Make label take full width */
        margin-bottom: 0.3rem;
    }
    .result-content {
        color: #333;
    }
    .correct-answer {
        background-color: #e6ffed; /* Light green */
        border-left-color: #4CAF50; /* Green */
    }
    .incorrect-answer {
        background-color: #ffeeee; /* Light red */
        border-left-color: #f44336; /* Red */
    }
    .user-answer-incorrect {
         color: #d32f2f; /* Red text for incorrect user answer */
         font-style: italic;
    }
    .correct-answer-text {
         color: #2e7d32; /* Green text for correct answer */
         font-weight: bold;
    }
    /* Sidebar styles */
    .stFileUploader label {
        font-size: 1.1rem; /* Increase font size for label */
        font-weight: bold;
        color: #333;
        margin-bottom: 0.5rem; /* Add space below label */
    }
    .sidebar .sidebar-content {
        background-color: #f8f9fa;
    }
    /* PDF Preview styling */
    .pdf-preview-container {
        margin-top: 1rem;
        border: 1px solid #ddd;
        border-radius: 5px;
        padding: 0.5rem;
    }
    .pdf-iframe {
        width: 100%;
        height: 400px; /* Adjust height as needed */
        border: none;
    }
    </style>
    """, unsafe_allow_html=True)

# Create PDF preview function
def display_pdf(file_bytes, file_name):
    st.markdown(f"##### Preview: {file_name}")
    try:
        base64_pdf = base64.b64encode(file_bytes).decode("utf-8")
        pdf_display = f'<div class="pdf-preview-container"><iframe class="pdf-iframe" src="data:application/pdf;base64,{base64_pdf}" type="application/pdf"></iframe></div>'
        st.markdown(pdf_display, unsafe_allow_html=True)
    except Exception as e:
        st.error(f"Error displaying PDF preview for {file_name}: {e}")

# Move file upload and preview to sidebar
with st.sidebar:
    #st.title("⚙️ Configuration")

    # Vector DB toggle
    st.header("🏫 Curriculum Database")
    use_vector_db = st.checkbox("Connect to my school's Curriculum Database", key="use_vector_db")
    vector_db_api_key = st.text_input(
        "LlamaIndex API Key (optional, overrides default)",
        type="password",
        value=VECTOR_DB_API_KEY or "",
        help="Enter your LlamaIndex API key if not set in environment/secrets."
    ) if use_vector_db else None
    st.markdown("---")
    # File uploader
    st.header("📄 Document Upload")
    uploaded_files = st.file_uploader(
        "Upload context documents (PDF, DOCX, PPTX)",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        help="Upload documents to provide context for flashcard generation."
    )

    process_files_button = st.button("Process Uploaded Files")

    if process_files_button and uploaded_files:
        st.session_state.source_content = "" # Reset content
        st.session_state.pdf_data = [] # Reset preview data
        st.session_state.pdf_names = [] # Reset names
        with st.spinner("Extracting text from documents..."):
            st.session_state.source_content = extract_text_from_multiple_files(uploaded_files)
            if st.session_state.source_content:
                 st.success(f"Successfully processed {len(uploaded_files)} document(s)!")
                 # Store file bytes and names for preview *after* successful processing
                 st.session_state.pdf_data = [f.getvalue() for f in uploaded_files if f.name.lower().endswith('.pdf')]
                 st.session_state.pdf_names = [f.name for f in uploaded_files if f.name.lower().endswith('.pdf')]
            else:
                 st.warning("Could not extract text from the uploaded file(s). Please check the files.")

    # Display PDF previews if data exists
    if st.session_state.pdf_data:
        st.markdown("---")
        st.subheader("PDF Previews")
        for i, pdf_bytes in enumerate(st.session_state.pdf_data):
             if i < len(st.session_state.pdf_names):
                 display_pdf(pdf_bytes, st.session_state.pdf_names[i])

    # Vector DB preview (if enabled and query present)
    if use_vector_db and st.session_state.get("query_input"):
        st.markdown("---")
        st.subheader("🔍 Retrieved Context Preview")
        with st.spinner("Retrieving relevant curriculum context..."):
            vector_context = retrieve_vector_context(
                st.session_state["query_input"],
                api_key=vector_db_api_key or VECTOR_DB_API_KEY
            )
            st.session_state["vector_context"] = vector_context
            if vector_context:
                st.markdown(f"<div style='font-size:0.95em; background:#f6f6fa; border-radius:8px; padding:0.7em; border:1px solid #eee; max-height:200px; overflow-y:auto; overflow-x:hidden;'>{vector_context[:1200]}{' ...' if len(vector_context)>1200 else ''}</div>", unsafe_allow_html=True)
            else:
                st.info("No relevant curriculum context found for this query.")

    # Question count configuration
    st.markdown("---")
    st.header("🔢 Question Counts")
    num_mc = st.number_input("Multiple Choice", min_value=0, max_value=7, value=3, key="num_mc")
    num_fb = st.number_input("Fill in the Blank", min_value=0, max_value=7, value=3, key="num_fb")
    num_tf = st.number_input("True/False", min_value=0, max_value=7, value=3, key="num_tf")
    num_qa = st.number_input("Question & Answer", min_value=0, max_value=7, value=3, key="num_qa")
    st.caption(f"Total questions: {num_mc + num_fb + num_tf + num_qa}")

# Main content area
st.title("֎ QuickQuizzer")
st.write("Enter a topic, use your school curriculum or upload files and generate practice questions or flashacards ✨")


# User input for the query
user_query = st.text_input("What do you want to practice today?", key="query_input")

# Place the two main action buttons side by side
col1, col2 = st.columns(2)

with col1:
    if st.button("🚀 Generate Practice Test", key="generate_button") and user_query:
        reset_flashcard_state()  # Clear flashcard mode state before generating practice test
        total_q = num_mc + num_fb + num_tf + num_qa
        if total_q == 0:
            st.warning("Please set the number of questions for at least one type in the sidebar.")
        else:
            loading_placeholder = st.empty()
            progress_bar = st.progress(0)
            
            # Show different message based on question count
            if total_q > 20:
                loading_placeholder.info(f"⏳ Generating {total_q} questions using optimized batch generation... this may take 30-60 seconds.")
            else:
                loading_placeholder.info("⏳ Generating practice test... this may take a moment.")

            # --- CONTEXT SELECTION LOGIC ---
            use_vector_db = st.session_state.get("use_vector_db", False)
            vector_context = st.session_state.get("vector_context", "")
            file_context = st.session_state.get("source_content", "")
            # Combine contexts if both are present
            if use_vector_db and vector_context and file_context:
                combined_context = vector_context + "\n\n" + file_context
            elif use_vector_db and vector_context:
                combined_context = vector_context
            elif file_context:
                combined_context = file_context
            else:
                combined_context = ""

            # Update progress as we generate
            progress_bar.progress(0.2)
            loading_placeholder.info(f"⏳ Generating questions in optimized batches... ({total_q} total)")
            
            flashcards_output = generate_flashcards(
                user_query,
                combined_context,
                int(num_mc), int(num_fb), int(num_tf), int(num_qa)
            )
            
            progress_bar.progress(1.0)

            if flashcards_output:
                progress_bar.empty()
                loading_placeholder.empty()
                # Parse and store questions and answers
                questions, answers, question_types, options_list, raw_answers = parse_questions(flashcards_output)

                # Enhanced validation with detailed feedback
                if not questions:
                    st.error("❌ No questions could be parsed from the LLM response.")
                    with st.expander("🔍 Debug Information"):
                        st.text("Raw LLM Response:")
                        st.text(flashcards_output[:3000] + "..." if len(flashcards_output) > 3000 else flashcards_output)
                    st.info("**Solutions:** Try regenerating with a clearer topic or fewer questions.")
                    
                elif len(questions) != total_q:
                    # Count by type for detailed report
                    type_counts = {'multiple_choice': 0, 'fill_blank': 0, 'true_false': 0, 'question_answer': 0, 'other': 0}
                    for q_type in question_types:
                        type_counts[q_type] = type_counts.get(q_type, 0) + 1
                    
                    # Calculate differences
                    expected_counts = {
                        'multiple_choice': num_mc,
                        'fill_blank': num_fb, 
                        'true_false': num_tf,
                        'question_answer': num_qa
                    }
                    
                    # Check if it's close enough (within 1-2 questions) to be usable
                    total_diff = abs(len(questions) - total_q)
                    max_type_diff = max(abs(type_counts[k] - expected_counts[k]) for k in expected_counts.keys())
                    
                    if total_diff <= 2 and max_type_diff <= 2:
                        # Close enough - show warning but allow usage
                        st.warning(f"⚠️ Got {len(questions)} questions instead of {total_q} (close enough to proceed)")
                        
                        # Show the differences
                        col1, col2 = st.columns(2)
                        with col1:
                            st.write("**Expected:**")
                            for q_type, count in expected_counts.items():
                                st.write(f"- {q_type.replace('_', ' ').title()}: {count}")
                        with col2:
                            st.write("**Actual:**")
                            for q_type, count in type_counts.items():
                                if count > 0:
                                    st.write(f"- {q_type.replace('_', ' ').title()}: {count}")
                    else:
                        # Too far off - show error
                        st.error(f"❌ Significant count mismatch: Expected {total_q} questions, got {len(questions)}")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            st.write("**Expected:**")
                            for q_type, count in expected_counts.items():
                                st.write(f"- {q_type.replace('_', ' ').title()}: {count}")
                        with col2:
                            st.write("**Actual:**")
                            for q_type, count in type_counts.items():
                                if count > 0:
                                    st.write(f"- {q_type.replace('_', ' ').title()}: {count}")
                        
                        with st.expander("🔍 Debug Information"):
                            st.text("Raw LLM Response (first 3000 characters):")
                            st.text(flashcards_output[:3000] + "..." if len(flashcards_output) > 3000 else flashcards_output)
                        
                        st.info("**Suggestions:**")
                        st.write("1. **Try Again** - The improved AI should be more consistent now")
                        st.write("2. **Reduce Numbers** - Try fewer questions per type (e.g., 2-5 each)")
                        st.write("3. **Simplify Topic** - Use more specific, concrete topics")
                        st.write("4. **Check Context** - Remove very long or complex uploaded documents")
                        
                        # Clear state and return early
                        for key in ['questions', 'answers', 'raw_answers', 'question_types', 'options_list']:
                            st.session_state[key] = []
                        st.stop()
                
                # Success case
                st.session_state.questions = questions
                st.session_state.answers = answers
                st.session_state.raw_answers = raw_answers
                st.session_state.question_types = question_types
                st.session_state.options_list = options_list
                # Reset state for the new session
                st.session_state.current_question_idx = 0
                st.session_state.user_answers = {}
                st.session_state.show_results = False
                
                if len(questions) == total_q:
                    loading_placeholder.success("✅ Practice Questions generated successfully!")
                else:
                    loading_placeholder.success(f"✅ Generated {len(questions)} questions (close to target {total_q})")
                time.sleep(1)
                loading_placeholder.empty()
                st.rerun()
            else:
                progress_bar.empty()
                loading_placeholder.error("Failed to generate flashcards. Please check the console for errors or try again.")

with col2:
    if st.button("📇 Generate Flash Cards", key="generate_flashcards_button") and user_query:
        reset_practice_test_state()  # Clear practice test mode state before generating flashcards
        with st.spinner("⏳ Generating 10 flashcards... this may take a moment."):
            # Only MCQs, 10 at a time - Using JSON mode for guaranteed valid output
            flashcard_system_prompt = (
                "You are an expert flashcard generator. You must respond with valid JSON only. "
                "Your response must be a JSON object with a 'flashcards' key containing an array of exactly 10 flashcards. "
                "Each flashcard must have exactly these 4 keys: 'question', 'options', 'correct', 'explanation'. "
                "The 'options' field must be an array of exactly 4 strings (no letter prefixes). "
                "The 'correct' field must be exactly 'A', 'B', 'C', or 'D' (capital letter only). "
                "Required JSON structure: "
                '{"flashcards": [{"question": "What is 2+2?", "options": ["3", "4", "5", "6"], "correct": "B", "explanation": "2+2 equals 4, which is option B."}, ...]}'
            )
            
            flashcard_user_prompt = f"Generate exactly 10 multiple-choice flashcards for the topic: {user_query}"
            
            # Use file/vector context if available
            use_vector_db = st.session_state.get("use_vector_db", False)
            vector_context = st.session_state.get("vector_context", "")
            file_context = st.session_state.get("source_content", "")
            combined_context = ""
            if use_vector_db and vector_context and file_context:
                combined_context = vector_context + "\n\n" + file_context
            elif use_vector_db and vector_context:
                combined_context = vector_context
            elif file_context:
                combined_context = file_context
            if combined_context:
                flashcard_user_prompt += f"\n\nUse this context to create relevant questions:\n{combined_context}"
            try:
                http_client = httpx.Client(timeout=60.0, trust_env=False)
                client = Groq(api_key=GROQ_API_KEY, http_client=http_client)
                completion = client.chat.completions.create(
                    model="llama3-70b-8192",
                    messages=[
                        {"role": "system", "content": flashcard_system_prompt},
                        {"role": "user", "content": flashcard_user_prompt}
                    ],
                    temperature=0.2,  # Lower temperature for consistency
                    top_p=0.9,
                    stream=False,
                    response_format={"type": "json_object"}  # JSON mode enabled!
                )
                
                # JSON mode guarantees valid JSON - no complex parsing needed!
                flashcards = json.loads(completion.choices[0].message.content)
                
                # Enhanced validation with detailed feedback
                valid_flashcards = []
                validation_errors = []
                
                # Handle both array and object responses (model might return object with array inside)
                if isinstance(flashcards, dict) and 'flashcards' in flashcards:
                    flashcards = flashcards['flashcards']
                elif isinstance(flashcards, dict) and len(flashcards) == 1:
                    # If dict with single key, extract the array
                    flashcards = list(flashcards.values())[0]
                
                if not isinstance(flashcards, list):
                    st.error("Invalid response format: Expected array of flashcards")
                    if st.button('🔄 Regenerate Flashcards', key='regenerate_flashcards_invalid_format'):
                        st.rerun()
                    st.stop()
                
                for i, card in enumerate(flashcards):
                    is_valid, error_msg = validate_flashcard_with_details(card, i)
                    if is_valid:
                        valid_flashcards.append(card)
                    else:
                        validation_errors.append(error_msg)
                
                # Show validation errors if any
                if validation_errors:
                    with st.expander("⚠️ Validation Issues Found", expanded=len(valid_flashcards) == 0):
                        for error in validation_errors:
                            st.warning(error)
                
                if not valid_flashcards:
                    st.error("No valid flashcards were generated. Please try again.")
                    if st.button('🔄 Try Again', key='regenerate_flashcards_no_valid'):
                        st.rerun()
                    st.stop()
                
                st.session_state.generated_flashcards = valid_flashcards
                st.session_state.flashcard_idx = 0
                
                success_msg = f"✅ Generated {len(valid_flashcards)} flashcards with JSON mode!"
                if len(validation_errors) > 0:
                    success_msg += f" ({len(validation_errors)} skipped due to validation errors)"
                
                st.success(success_msg)
                time.sleep(1)
                st.rerun()
                        
            except Exception as e:
                st.error(f"Failed to generate flashcards: {e}")
                if st.button('🔄 Try Again', key='regenerate_flashcards_exception'):
                    st.rerun()

# Function to display results
def generate_all_explanations():
    """Generate LLM explanations for all questions and store in session state."""
    if 'answer_explanations' not in st.session_state:
        st.session_state.answer_explanations = {}
    if 'answer_scores' not in st.session_state:
        st.session_state.answer_scores = {}
    total = len(st.session_state.questions)
    for idx in range(total):
        q_type = st.session_state.question_types[idx]
        question_html = st.session_state.questions[idx]
        correct_answer_raw = st.session_state.raw_answers[idx]
        user_answer = st.session_state.user_answers.get(idx, None)
        # Only generate if not already present
        if idx not in st.session_state.answer_explanations:
            if user_answer is not None:
                score, explanation = evaluate_answer_llm(
                    question_html, correct_answer_raw, user_answer, q_type
                )
            else:
                score, explanation = 0.0, "No answer provided."
            st.session_state.answer_explanations[idx] = explanation
            st.session_state.answer_scores[idx] = score

def generate_session_review_llm(api_key=GROQ_API_KEY):
    """
    Use the LLM to generate an overall review/observation of the session for the user.
    Returns a review string.
    """
    questions = st.session_state.questions
    question_types = st.session_state.question_types
    correct_answers = st.session_state.raw_answers
    user_answers = [st.session_state.user_answers.get(i, None) for i in range(len(questions))]
    explanations = [st.session_state.answer_explanations.get(i, "") for i in range(len(questions))]
    scores = [st.session_state.answer_scores.get(i, 0.0) for i in range(len(questions))]

    # Prepare a compact but informative session summary
    session_data = []
    for i in range(len(questions)):
        session_data.append({
            "question_type": question_types[i],
            "question": questions[i],
            "correct_answer": correct_answers[i],
            "user_answer": user_answers[i],
            "score": scores[i],
            "explanation": explanations[i],
        })

    # Calculate summary statistics
    total_score = sum(scores)
    average_score = total_score / len(scores) if scores else 0
    
    system_prompt = (
        "You are an expert tutor reviewing a student's practice session. "
        "Generate a structured JSON review with specific sections. "
        "Be encouraging, constructive, and specific in your feedback. "
        "Respond with JSON in this exact format: "
        '{"overall_performance": "summary text", "strengths": ["strength1", "strength2"], "areas_for_improvement": ["area1", "area2"], "actionable_advice": ["tip1", "tip2"], "motivational_message": "encouraging text"}'
    )
    
    user_prompt = (
        f"Analyze this practice session with {len(questions)} questions.\n"
        f"Average score: {average_score:.2f}/1.0\n"
        f"Total score: {total_score:.2f}/{len(questions)}\n\n"
        f"Session details:\n{json.dumps(session_data, indent=2)}\n\n"
        f"Generate a structured review focusing on question types, performance patterns, and specific study recommendations."
    )
    
    try:
        http_client = httpx.Client(timeout=90.0, trust_env=False)
        client = Groq(api_key=api_key, http_client=http_client)
        completion = client.chat.completions.create(
            model="llama3-70b-8192",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=0.2,
            top_p=0.95,
            stream=False,
            response_format={"type": "json_object"}  # JSON mode enabled!
        )
        
        # JSON mode guarantees valid JSON - parse and format nicely
        review_data = json.loads(completion.choices[0].message.content)
        
        # Format the structured review into readable text
        formatted_review = f"""**Overall Performance:** {review_data.get('overall_performance', 'Good effort on this session!')}

**Your Strengths:**
{chr(10).join(f"• {strength}" for strength in review_data.get('strengths', ['You completed the session!']))}

**Areas for Improvement:**
{chr(10).join(f"• {area}" for area in review_data.get('areas_for_improvement', ['Keep practicing!']))}

**Actionable Study Tips:**
{chr(10).join(f"• {tip}" for tip in review_data.get('actionable_advice', ['Review the material again']))}

**Keep Going!** {review_data.get('motivational_message', 'You are making progress - keep up the great work!')}"""
        
        return formatted_review
        
    except Exception as e:
        return f"(Could not generate review: {e})"

def display_results():
    st.header("\U0001F3C1 Session Results")
    st.write("Here's how you did:")

    # Generate explanations if not already done
    if 'answer_explanations' not in st.session_state or len(st.session_state.answer_explanations) < len(st.session_state.questions):
        with st.spinner("Generating explanations for your answers..."):
            generate_all_explanations()

    # Generate session review if not already done
    if 'session_review' not in st.session_state:
        with st.spinner("Generating overall review of your session..."):
            st.session_state.session_review = generate_session_review_llm()

    # Display the review at the bottom with enhanced UI
    #st.markdown("---")
    # Clean up LLM review output to remove redundant headers
    review_text = st.session_state.session_review
    # Remove leading markdown headers like '**Session Review**', '# Session Review', etc.
    review_text = re.sub(r"^(\s*(\*\*|#)+\s*Session Review\s*(\*\*)?\s*:?\s*)", "", review_text, flags=re.IGNORECASE)
    st.markdown(
        f'''<div style="background: linear-gradient(90deg, #e0e7ff 0%, #f0f4ff 100%); border-radius: 16px; padding: 2em 1.5em; border: 2px solid #4a90e2; font-size:1.15em; margin-top: 2em; margin-bottom: 2em; box-shadow: 0 4px 16px rgba(74,144,226,0.08);">
        <span style="font-size:1.5em; vertical-align:middle; margin-right:0.5em;">📝</span>
        <span style="font-weight:600; color:#2d3a5a;">Overall Session Review</span>
        <hr style="border: none; border-top: 1.5px solid #b3c6ff; margin: 0.7em 0 1.2em 0;">
        <div style="color:#222; line-height:1.7;">{review_text}</div>
        </div>''',
        unsafe_allow_html=True
    )
    #st.markdown("---")

    score = 0
    total = len(st.session_state.questions)
    qa_total_score = 0.0  # For Q&A LLM-based scoring
    qa_count = 0

    for idx in range(total):
        q_type = st.session_state.question_types[idx]
        question_html = st.session_state.questions[idx]
        correct_answer_display = st.session_state.answers[idx] # Formatted HTML answer
        correct_answer_raw = st.session_state.raw_answers[idx] # Raw answer string
        user_answer = st.session_state.user_answers.get(idx, None) # Get user's stored answer

        is_correct = False
        user_answer_display = "Not Answered" if user_answer is None else str(user_answer)
        llm_score = st.session_state.answer_scores.get(idx, 0.0)
        llm_explanation = st.session_state.answer_explanations.get(idx, "")

        # Determine Correctness
        if user_answer is not None:
            if q_type == "multiple_choice":
                correct_option_prefix = correct_answer_raw.strip()[0].upper() # Get 'A' from 'A)' or 'A'
                is_correct = str(user_answer).strip().upper() == correct_option_prefix
                options = st.session_state.options_list[idx]
                for opt in options:
                     if opt.strip().startswith(str(user_answer).strip().upper()):
                         user_answer_display = opt
                         break
            elif q_type == "true_false":
                is_correct = str(user_answer).strip().lower() == correct_answer_raw.strip().lower()
            elif q_type == "fill_blank":
                is_correct = str(user_answer).strip().lower() == correct_answer_raw.strip().lower()
            elif q_type == "question_answer":
                qa_total_score += llm_score
                qa_count += 1
                is_correct = llm_score >= 0.7  # Consider correct if score >= 0.7
            else:
                 is_correct = str(user_answer).strip().lower() == correct_answer_raw.strip().lower()

            if is_correct is True and q_type != "question_answer":
                score += 1

        # Display Result Item
        st.markdown(f"---") # Separator
        st.markdown(f'<div class="result-header">Question {idx + 1} ({q_type.replace("_", " ").title()})</div>', unsafe_allow_html=True)

        # Display Question Content
        st.markdown(f'<div class="result-section"><span class="result-label">Question:</span><div class="result-content">{question_html}</div></div>', unsafe_allow_html=True)

        # Display User's Answer
        user_answer_class = ""
        if is_correct is False:
             user_answer_class = "user-answer-incorrect"
        elif is_correct is True:
             user_answer_class = "correct-answer-text" # Or keep default

        st.markdown(f'<div class="result-section { "incorrect-answer" if is_correct is False else "" }"><span class="result-label">Your Answer:</span><div class="result-content {user_answer_class}">{user_answer_display}</div></div>', unsafe_allow_html=True)

        # Display Individual Score
        indiv_score = llm_score if user_answer is not None else 0
        st.markdown(f'<div class="result-section"><span class="result-label">Score:</span> <span class="result-content">{indiv_score:.2f} / 1.00</span></div>', unsafe_allow_html=True)

        # Display Correct Answer (always show)
        if q_type == "multiple_choice":
            options_html = '<ul style="list-style-type: none; padding-left: 0;">'
            for opt in st.session_state.options_list[idx]:
                 options_html += f'<li>{opt}</li>'
            options_html += '</ul>'
            st.markdown(f'<div class="result-section"><span class="result-label">Options:</span><div class="result-content">{options_html}</div></div>', unsafe_allow_html=True)

        st.markdown(f'<div class="result-section correct-answer"><div class="result-content correct-answer-text">{correct_answer_display}</div></div>', unsafe_allow_html=True)

        # Show LLM explanation for all question types
        st.markdown(f'<div class="result-section"><span class="result-label">Explanation:</span> <span class="result-content">{llm_explanation}</span></div>', unsafe_allow_html=True)

    # Final Score
    st.markdown("---")
    st.header("\U0001F4CA Final Score")
    # Add LLM Q&A score to total
    total_non_qa = total - qa_count
    total_score = score + qa_total_score
    total_possible = total_non_qa + qa_count  # Each Q&A max 1 point
    if total_possible > 0:
        percentage = (total_score / total_possible) * 100
        st.metric(label="Accuracy", value=f"{total_score:.2f}/{total_possible}", delta=f"{percentage:.1f}%")
    else:
        st.write("No questions were generated or answered.")

    # Reset Button
    if st.button("\U0001F504 Start New Session", key="reset_button"):
        keys_to_clear = [
            'questions', 'answers', 'raw_answers', 'question_types', 'options_list',
            'current_question_idx', 'user_answers', 'show_results', 'source_content',
            'pdf_data', 'pdf_names', 'answer_explanations', 'answer_scores', 'session_review'
        ]
        for key in keys_to_clear:
            if key in st.session_state:
                del st.session_state[key]
        st.session_state.current_question_idx = 0
        st.session_state.user_answers = {}
        st.session_state.show_results = False
        st.session_state.source_content = ""
        st.session_state.pdf_data = []
        st.session_state.pdf_names = []
        st.rerun()

# Main Logic: Display Questions or Results
if st.session_state.get('show_results', False):
    display_results()

elif st.session_state.get('questions'):
    # Display Current Question
    current_idx = st.session_state.current_question_idx
    total_questions = len(st.session_state.questions)

    # Progress bar
    progress = (current_idx + 1) / total_questions
    st.progress(progress)
    st.markdown(f'<div class="question-counter">Question {current_idx + 1} of {total_questions}</div>', unsafe_allow_html=True)

    # Get current question details
    current_question = st.session_state.questions[current_idx]
    current_type = st.session_state.question_types[current_idx]
    current_options = st.session_state.options_list[current_idx]

    # Display question in a styled container
    st.markdown(f'<div class="flashcard">{current_question}</div>', unsafe_allow_html=True)

    # Input Widget based on Question Type
    user_input = None
    key_prefix = f"q_{current_idx}_{current_type}"

    if current_type == "multiple_choice":
        st.write("**Select your answer:**")
        mcq_options = [opt for opt in current_options] # Ensures we have a list
        # Display options as radio buttons
        user_input = st.radio(
             "Options",
             mcq_options,
             key=f"{key_prefix}_mcq",
             index=None, # Default to no selection
             label_visibility="collapsed" # Hide the 'Options' label itself
        )
        if user_input:
             st.session_state.user_answers[current_idx] = user_input.split(')')[0] # Store 'A', 'B' etc.

    elif current_type == "true_false":
        st.write("**Select True or False:**")
        user_input = st.radio(
            "Answer",
            ["True", "False"],
            key=f"{key_prefix}_tf",
            index=None, # Default to no selection
             label_visibility="collapsed"
        )
        if user_input:
            st.session_state.user_answers[current_idx] = user_input

    elif current_type == "fill_blank":
        st.write("**Fill in the blank:**")
        user_input = st.text_input(
            "Your answer",
            key=f"{key_prefix}_fb",
            value=st.session_state.user_answers.get(current_idx, "") # Keep previous input if navigating back
        )
        st.session_state.user_answers[current_idx] = user_input # Store continuously

    elif current_type == "question_answer":
        st.write("**Your Answer:**")
        user_input = st.text_area(
            "Provide your detailed answer",
            key=f"{key_prefix}_qa",
            height=150,
            value=st.session_state.user_answers.get(current_idx, "") # Keep previous input
        )
        st.session_state.user_answers[current_idx] = user_input # Store continuously

    else: # Default/Other types - Treat as text input
        st.write("**Your Answer:**")
        user_input = st.text_input(
            "Your answer",
            key=f"{key_prefix}_other",
            value=st.session_state.user_answers.get(current_idx, "")
        )
        st.session_state.user_answers[current_idx] = user_input

    # Navigation Buttons
    st.markdown('<div class="nav-buttons">', unsafe_allow_html=True) # Start flex container

    # 'Previous' Button Column
    prev_col, counter_col, next_col = st.columns([1, 2, 1])

    with prev_col:
        if current_idx > 0:
            if st.button("⬅️ Previous", key="prev_button"):
                st.session_state.current_question_idx -= 1
                st.rerun()
        else:
            st.write("") # Placeholder to maintain layout

    # Placeholder for counter (already shown above)
    with counter_col:
         st.write("") # Keep space

    # 'Next' / 'Finish' Button Column
    with next_col:
        is_last_question = current_idx == total_questions - 1
        button_text = "🏁 Finish Session" if is_last_question else "Next ➡️"
        button_key = "finish_button" if is_last_question else "next_button"

        if st.button(button_text, key=button_key):
            if is_last_question:
                # Check if all questions have been attempted (optional, but good UX)
                if len(st.session_state.user_answers) < total_questions:
                     st.warning("You haven't answered all questions. Are you sure you want to finish?")
                     # Could add another confirmation button here if needed
                st.session_state.show_results = True
            else:
                st.session_state.current_question_idx += 1
            st.rerun()

    st.markdown('</div>', unsafe_allow_html=True) # End flex container

# --- ANIMATED HTML/CSS/JS FLIP CARD PROTOTYPE ---
def get_flipcard_html(question, options_html, explanation, correct):
    return f'''
<style>
.flip-card-container {{
  display: flex;
  justify-content: center;
  align-items: center;
  min-height: 400px;
  margin-bottom: 2.5em;
}}
.flip-card {{
  background: linear-gradient(135deg, #e0e7ff 0%, #f0f4ff 100%);
  width: 540px;
  min-height: 340px;
  perspective: 1200px;
  cursor: pointer;
  border-radius: 28px;
  box-shadow: 0 8px 32px rgba(74,144,226,0.13), 0 1.5px 8px rgba(0,0,0,0.07);
  border: 2.5px solid #b3c6ff;
  box-sizing: border-box;
  position: relative;
  transition: height 0.3s cubic-bezier(.4,2,.6,1);
}}
.flip-card-inner {{
  position: relative;
  width: 100%;
  transition: transform 0.8s cubic-bezier(.4,2,.6,1), height 0.3s cubic-bezier(.4,2,.6,1);
  transform-style: preserve-3d;
  box-sizing: border-box;
}}
.flip-card.flipped .flip-card-inner {{
  transform: rotateY(180deg);
}}
.flip-card-front, .flip-card-back {{
  position: absolute;
  top: 0;
  left: 0;
  width: 100%;
  height: 100%;
  backface-visibility: hidden;
  border-radius: 28px;
  box-shadow: none;
  background: transparent;
  padding: 2.2em;
  display: flex;
  flex-direction: column;
  justify-content: flex-start;
  align-items: stretch;
  box-sizing: border-box;
}}
.question-text {{
  font-weight: 700;
  font-size: 1.18em;
  margin-bottom: 1.2em;
  width: 100%;
  box-sizing: border-box;
  padding: 0.6em 0.7em;
  word-break: break-word;
  overflow-wrap: break-word;
  line-height: 1.5;
  background: rgba(255,255,255,0.85);
  border-radius: 14px;
  box-shadow: 0 1.5px 8px rgba(0,0,0,0.04);
  margin-top: 0;
}}
.options-grid {{
  display: grid;
  grid-template-columns: 1fr 1fr;
  gap: 1.1em 1.2em;
  margin: 0;
  width: 100%;
  justify-items: stretch;
  box-sizing: border-box;
}}
.flip-card-front {{
  z-index: 2;
}}
.flip-card-back {{
  transform: rotateY(180deg);
  background: linear-gradient(135deg, #f6f8ff 0%, #e0e7ff 100%);
  z-index: 3;
  box-shadow: none;
}}
.option-btn {{
  display: block;
  width: 100%;
  margin: 0;
  padding: 1em 0.8em;
  border: none;
  border-radius: 11px;
  font-size: 1.05em;
  font-weight: 600;
  background: #f3f7fa;
  color: #1a237e;
  cursor: pointer;
  transition: background 0.2s, color 0.2s, box-shadow 0.2s;
  text-align: left;
  box-shadow: 0 1.5px 8px rgba(0,0,0,0.04);
  outline: none;
  box-sizing: border-box;
}}
.option-btn.selected.correct {{
  background: linear-gradient(90deg, #b2f7ef 0%, #a8ff78 100%) !important;
  color: #00695c !important;
  box-shadow: 0 0 0 2px #00b89433;
}}
.option-btn.selected.incorrect {{
  background: linear-gradient(90deg, #ffd6e0 0%, #ffb199 100%) !important;
  color: #b71c1c !important;
  box-shadow: 0 0 0 2px #ff767533;
}}
.option-btn.correct {{
  background: linear-gradient(90deg, #b2f7ef 0%, #a8ff78 100%) !important;
  color: #00695c !important;
  box-shadow: 0 0 0 2px #00b89433;
}}
.option-btn:disabled {{
  opacity: 0.92;
  cursor: not-allowed;
}}
.flip-label {{
  font-size: 1.05em;
  color: #6c7a89;
  margin-top: 1.2em;
  text-align: center;
  width: 100%;
  letter-spacing: 0.01em;
}}
@media (max-width: 700px) {{
  .flip-card {{
    width: 98vw;
    min-width: 0;
    max-width: 99vw;
  }}
}}
</style>
<div class="flip-card-container">
  <div class="flip-card" id="flipCard">
    <div class="flip-card-inner" id="flipCardInner">
      <div class="flip-card-front" id="flipCardFront">
        <div class="question-text">Q: {question}</div>
        <div id="optionsArea" class="options-grid">
          {options_html}
        </div>
        <div class="flip-label">Click the card (not options) to flip for explanation</div>
      </div>
      <div class="flip-card-back" id="flipCardBack">
        <div style="font-weight:700;font-size:1.15em;margin-bottom:1.2em;">Explanation</div>
        <div style="font-size:1.05em;line-height:1.6;">{explanation}</div>
        <div class="flip-label">Click the card to flip back</div>
      </div>
    </div>
  </div>
</div>
<script>
var correctOption = "{correct}";
function setFlipCardHeight() {{
  var front = document.getElementById('flipCardFront');
  var card = document.getElementById('flipCard');
  var inner = document.getElementById('flipCardInner');
  if (front && card && inner) {{
    card.style.height = 'auto';
    inner.style.height = 'auto';
    var frontHeight = front.offsetHeight;
    card.style.height = frontHeight + 'px';
    inner.style.height = frontHeight + 'px';
  }}
}}
window.addEventListener('load', setFlipCardHeight);
window.addEventListener('resize', setFlipCardHeight);
document.getElementById('flipCard').addEventListener('click', function(e) {{
  if (!e.target.classList.contains('option-btn')) {{
    this.classList.toggle('flipped');
    setTimeout(setFlipCardHeight, 350); // after flip animation
  }}
}});
function selectOption(label) {{
  // Remove previous selection classes
  var btns = document.querySelectorAll('.option-btn');
  btns.forEach(function(btn) {{
    btn.classList.remove('selected', 'correct', 'incorrect');
  }});
  // Find the correct and selected buttons
  var selectedBtn = null;
  var correctBtn = null;
  btns.forEach(function(btn) {{
    if (btn.textContent.trim().startsWith(label + ')')) {{
      selectedBtn = btn;
    }}
    if (btn.textContent.trim().startsWith(correctOption + ')')) {{
      correctBtn = btn;
    }}
  }});
  if (selectedBtn) {{
    if (label === correctOption) {{
      selectedBtn.classList.add('selected', 'correct');
    }} else {{
      selectedBtn.classList.add('selected', 'incorrect');
      if (correctBtn) correctBtn.classList.add('correct');
    }}
  }}
}}
</script>
'''

if st.session_state.get('generated_flashcards'):
    st.markdown('---')
    st.markdown("<h2 style='text-align:center;'>Flashcards</h2>", unsafe_allow_html=True)
    flashcards = st.session_state.generated_flashcards
    idx = st.session_state.get('flashcard_idx', 0)
    total = len(flashcards)
    card = flashcards[idx]
    if not is_valid_flashcard(card):
        st.error(f"Malformed flashcard at index {idx}: {card}")
        st.stop()
    # Prepare data for JS
    question = card['question']
    options = card['options']
    correct = card['correct']
    explanation = card['explanation']
    # Build options HTML as a grid
    options_html = ""
    for i, opt in enumerate(options):
        label = chr(65 + i)
        # Add event.stopPropagation() to prevent card flip on option click
        options_html += f'<button class="option-btn" id="opt_{i}" onclick="event.stopPropagation();selectOption(\'{label}\')" type="button">{label}) {opt}</button>'
    html_code = get_flipcard_html(question, options_html, explanation, correct)
    components.html(html_code, height=420)
    # Navigation below the card
    nav_cols = st.columns([1, 2, 1])
    with nav_cols[0]:
        if idx > 0:
            if st.button("⬅️ Previous Card", key="prev_flashcard_animated"):
                st.session_state.flashcard_idx = idx - 1
                st.rerun()
    with nav_cols[2]:
        if idx < total - 1:
            if st.button("Next Card ➡️", key="next_flashcard_animated"):
                st.session_state.flashcard_idx = idx + 1
                st.rerun()
